"""
PPTX Renderer (Step 3)

Generates a competency matrix deck using the "Overlay" master template technique:
- Table/borders are static in the slide master
- Named placeholders (Selection Pane) overlay the table cells

This module is intentionally deterministic and LLM-free.
"""

from __future__ import annotations

import datetime as _dt
import re
from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Any, Dict, Iterable, List, Mapping, Optional, Tuple, Union

from pptx import Presentation

from ..config import (
    GCS_BUCKET_NAME,
    GCS_PREFIX,
    GCS_RETURN_SIGNED_URL,
    GCS_SIGNED_URL_TTL_SECONDS,
    TEMPLATE_PPTX,
)


REQUIRED_PLACEHOLDERS = [
    # Metadata overlay
    "Header_CompType",
    "Header_JobGroup",
    "Header_Competency",
    "Header_Department",
    # Content
    "Main_Title",
    "Main_Definition",
    "Topic_Title",
    "Topic_Description",
    # Columns
    "Level_Expert",
    "Level_Advanced",
    "Level_Intermediate",
    "Level_Beginner",
]


def _as_dict(obj: Any) -> Dict[str, Any]:
    """Support pydantic models or plain dicts."""
    if obj is None:
        return {}
    if isinstance(obj, dict):
        return obj
    # pydantic v2
    if hasattr(obj, "model_dump"):
        return obj.model_dump()
    # pydantic v1
    if hasattr(obj, "dict"):
        return obj.dict()
    raise TypeError(f"Unsupported object type for jobs_data item: {type(obj)!r}")


def _normalize_lines(text_or_list: Any) -> str:
    """Normalize text into newline-separated items WITHOUT explicit bullet glyphs.

    The master template placeholders are expected to already have bullet/list
    formatting. Therefore we remove leading bullet symbols (•, -, *, etc.) and
    keep one item per line.
    """

    if text_or_list is None:
        return ""

    if isinstance(text_or_list, (list, tuple)):
        items = [str(x).strip() for x in text_or_list if str(x).strip()]
    else:
        s = str(text_or_list).strip()
        if not s:
            return ""
        items = [ln.strip() for ln in s.splitlines() if ln.strip()]

    cleaned: List[str] = []
    for it in items:
        # Strip common leading bullet markers
        it2 = re.sub(r"^\s*[•\-\*•]\s*", "", it).strip()
        # Remove stray bullet glyphs that appear in the middle
        it2 = it2.replace("•", "").strip()
        if it2:
            cleaned.append(it2)

    return "\n".join(cleaned)


def _strip_bullet_symbols(text: str) -> str:
    """Remove explicit bullet glyphs from any string before placing in PPTX."""

    if not text:
        return text
    return str(text).replace("•", "").strip()


def _find_layout(prs: Presentation, layout_name: str) -> Any:
    for layout in prs.slide_layouts:
        if getattr(layout, "name", None) == layout_name:
            return layout
    # fallback: last layout
    return prs.slide_layouts[-1]


def _layout_placeholder_map(layout: Any) -> Dict[str, int]:
    """
    Returns: {placeholder_name: placeholder_idx}
    placeholder_idx is slide placeholder index (placeholder_format.idx)
    """
    m: Dict[str, int] = {}
    for shape in layout.placeholders:
        try:
            m[shape.name] = shape.placeholder_format.idx
        except Exception:
            # Not a placeholder or missing idx; ignore
            continue
    return m


def validate_template_placeholders(template_path: Union[str, Path], layout_name: str = "Competency_Layout") -> Tuple[bool, List[str]]:
    """
    Checks whether the template contains the named placeholders required by the overlay mapping.
    Returns (ok, missing_names).
    """
    prs = Presentation(str(template_path))
    layout = _find_layout(prs, layout_name)
    m = _layout_placeholder_map(layout)
    missing = [n for n in REQUIRED_PLACEHOLDERS if n not in m]
    return (len(missing) == 0, missing)


def generate_competency_slides(
    template_path: Union[str, Path],
    output_path: Union[str, Path],
    jobs_data: Iterable[Any],
    *,
    layout_name: str = "Competency_Layout",
    strict: bool = True,
) -> Path:
    """
    Generate a PPTX where:
      - one job (single element in jobs_data) can have 2-3 topics
      - one topic == one slide

    Args:
      template_path: path to master template.pptx
      output_path: output pptx path
      jobs_data: list of job dicts (or pydantic objects) matching Step-2 JSON contract
      layout_name: expected slide layout name inside template
      strict: if True, raises when required placeholders are missing

    Returns:
      output_path as Path
    """
    template_path = Path(template_path)
    output_path = Path(output_path)

    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    prs = Presentation(str(template_path))
    layout = _find_layout(prs, layout_name)
    layout_map = _layout_placeholder_map(layout)

    if strict:
        missing = [n for n in REQUIRED_PLACEHOLDERS if n not in layout_map]
        if missing:
            raise ValueError(
                "Template layout is missing required placeholders: "
                + ", ".join(missing)
                + f" (layout_name={layout_name})"
            )

    jobs_list = [_as_dict(j) for j in jobs_data]
    if not jobs_list:
        raise ValueError("jobs_data is empty. Expected one job with 2-3 topics.")

    # One job per run (enforce; but allow list length > 1 if caller passes multiple)
    for job in jobs_list:
        topics = job.get("topics") or []
        if not isinstance(topics, list) or len(topics) == 0:
            raise ValueError("Job has no topics. Step-2 must generate 2-3 topics.")

        for topic in topics:
            slide = prs.slides.add_slide(layout)

            # Normalize topic proficiency fields
            expert = _normalize_lines(topic.get("expert"))
            advanced = _normalize_lines(topic.get("advanced"))
            intermediate = _normalize_lines(topic.get("intermediate"))
            beginner = _normalize_lines(topic.get("beginner"))

            content_map: Mapping[str, Any] = {
                # Metadata Overlay
                "Header_CompType": job.get("general_group") or job.get("comp_type", ""),
                "Header_JobGroup": job.get("specific_group") or job.get("job_group", ""),
                "Header_Competency": job.get("competency_name", ""),
                "Header_Department": job.get("job_location") or job.get("department", ""),

                # Content
                "Main_Title": job.get("competency_name", ""),
                "Main_Definition": job.get("definition", ""),
                "Topic_Title": topic.get("title", ""),
                "Topic_Description": topic.get("desc", ""),

                # Columns
                "Level_Expert": expert,
                "Level_Advanced": advanced,
                "Level_Intermediate": intermediate,
                "Level_Beginner": beginner,
            }

            # Write to placeholders
            for name, value in content_map.items():
                if name not in layout_map:
                    # Template mismatch; skip in non-strict mode
                    if strict:
                        raise KeyError(f"Placeholder not found in template layout: {name}")
                    continue

                ph_idx = layout_map[name]
                try:
                    slide.placeholders[ph_idx].text = "" if value is None else _strip_bullet_symbols(str(value))
                except KeyError:
                    # Some templates may not carry placeholders into slides; fail only if strict
                    if strict:
                        raise KeyError(f"Slide missing placeholder idx={ph_idx} for '{name}'")
                except Exception as e:
                    if strict:
                        raise RuntimeError(f"Failed writing to placeholder '{name}' (idx={ph_idx}): {e}") from e

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def default_output_name(job_title: str, *, when: Optional[_dt.date] = None) -> str:
    """
    Utility: build {job_title}_{date}.pptx with safe filename characters.
    Forces underscores instead of spaces to ensure clickable Markdown links.
    """
    when = when or _dt.date.today()
    # Allow Arabic chars, English chars, numbers, and hyphens. Remove everything else.
    safe = re.sub(r"[^\w\u0600-\u06FF\- ]+", "", job_title, flags=re.UNICODE).strip()
    # CRITICAL FIX: Replace ALL whitespace (spaces, tabs) with underscores
    safe = re.sub(r"\s+", "_", safe) or "job"
    return f"{safe}_{when.isoformat()}.pptx"


def _get_storage_client():
    """Create a GCS client using Application Default Credentials."""
    from google.cloud import storage

    return storage.Client()


def _join_blob_name(prefix: str, filename: str) -> str:
    prefix = (prefix or "").strip().strip("/")
    return f"{prefix}/{filename}" if prefix else filename


def _upload_pptx_to_gcs(local_path: Path, *, bucket_name: str, blob_name: str) -> str:
    """Upload PPTX to GCS and return a console-style HTTPS URL.

    The returned URL matches the format you requested:
      https://storage.cloud.google.com/BUCKET/OBJECT

    Note: This URL often requires the viewer to have GCP permissions. If you
    need link-based access without auth, enable signed URLs.
    """

    client = _get_storage_client()
    bucket = client.bucket(bucket_name)
    blob = bucket.blob(blob_name)
    blob.upload_from_filename(
        str(local_path),
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    return f"https://storage.cloud.google.com/{bucket.name}/{blob.name}"


def _maybe_generate_signed_url(*, bucket_name: str, blob_name: str, ttl_seconds: int) -> str | None:
    """Return a V4 signed download URL if possible; otherwise None."""

    if ttl_seconds <= 0:
        return None
    try:
        import datetime as dt
        from google.cloud import storage

        client = storage.Client()
        bucket = client.bucket(bucket_name)
        blob = bucket.blob(blob_name)
        return blob.generate_signed_url(
            version="v4",
            expiration=dt.timedelta(seconds=ttl_seconds),
            method="GET",
        )
    except Exception:
        return None


# -----------------------------
# ADK tool wrapper (Step 3 + Step 4)
# -----------------------------

import base64

async def render_competency_pptx(
    jobs_data: List[Dict[str, Any]],
    job_title: str = "job",
    output_filename: Optional[str] = None,
    template_artifact_filename: Optional[str] = None,
    layout_name: str = "Competency_Layout",
    strict: bool = True,
    tool_context=None,
) -> Dict[str, Any]:
    """
    [Step 3] Generates the final PowerPoint file and uploads it to Cloud Storage.

    This tool takes the structured JSON data (from Step 2) and overlays it onto
    the official government PPTX template.

    Args:
        jobs_data: List of competency dictionaries (output from generate_competency_model).
        job_title: The name of the job (used for the filename).
        output_filename: Optional override for the .pptx filename.
        template_artifact_filename: Optional custom template (defaults to system template).
        layout_name: The slide master layout to use (default: 'Competency_Layout').
        strict: If True, fails if template placeholders are missing.
        tool_context: ADK context (injected automatically).

    Returns:
        A dictionary containing the 'final_message' (in Arabic) with the download link,
        upload status, and GCS details.
    """

    from google.genai import types

    out_name = output_filename or default_output_name(job_title)
    topics_count = 0
    if jobs_data and isinstance(jobs_data, list):
        topics_count = sum(len((j or {}).get("topics") or []) for j in jobs_data)

    bucket = GCS_BUCKET_NAME
    prefix = GCS_PREFIX
    signed = GCS_RETURN_SIGNED_URL
    ttl = GCS_SIGNED_URL_TTL_SECONDS

    with TemporaryDirectory() as td:
        td_path = Path(td)
        out_path = td_path / out_name

        # Resolve template path
        template_path = Path(TEMPLATE_PPTX)
        if template_artifact_filename and tool_context is not None:
            try:
                part = await tool_context.load_artifact(filename=template_artifact_filename)
                data_bytes = None
                if isinstance(part, dict) and "inlineData" in part:
                    b64 = part["inlineData"].get("data")
                    if b64: data_bytes = base64.b64decode(b64)
                elif part and getattr(part, "inline_data", None):
                    data_bytes = part.inline_data.data
                elif part and hasattr(part, "data"):
                    data_bytes = part.data
                
                if data_bytes:
                    template_path = td_path / "template.pptx"
                    template_path.write_bytes(data_bytes)
            except Exception:
                template_path = Path(TEMPLATE_PPTX)

        # Step 3: render
        generate_competency_slides(
            template_path=template_path,
            output_path=out_path,
            jobs_data=jobs_data,
            layout_name=layout_name,
            strict=strict,
        )

        # Step 4: upload to GCS (default)
        blob_name = _join_blob_name(prefix, out_name)
        gcs_url: Optional[str] = None
        signed_url: Optional[str] = None
        upload_error: Optional[str] = None

        if bucket:
            try:
                gcs_url = _upload_pptx_to_gcs(out_path, bucket_name=bucket, blob_name=blob_name)
                if signed:
                    signed_url = _maybe_generate_signed_url(bucket_name=bucket, blob_name=blob_name, ttl_seconds=ttl)
            except Exception as e:
                upload_error = str(e)
        else:
            upload_error = "GCS_BUCKET_NAME is not set"

        pptx_bytes = out_path.read_bytes()

        artifact_saved = False
        artifact_filename = None
        if gcs_url:
            final_message = (
                "تم إنشاء مصفوفة الكفاءات بنجاح\n\n"
                f"يمكنك الوصول إلى المستند من هنا: [{out_name}]({gcs_url})"
            )
            if signed_url:
                final_message += (
                    "\n\n(رابط موقع): "
                    f"[{out_name}]({signed_url})"
                )
        else:
            # Optional artifact fallback (useful in local dev / if upload fails)
            artifact_saved = False
            artifact_filename = None
            if tool_context is not None:
                try:
                    await tool_context.save_artifact(
                        filename=out_name,
                        artifact=types.Part.from_bytes(
                            data=pptx_bytes,
                            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        ),
                    )
                    artifact_saved = True
                    artifact_filename = out_name
                except Exception:
                    artifact_saved = False
            final_message = (
                "تم إنشاء مصفوفة الكفاءات بنجاح، ولكن تعذر الرفع إلى السحابة.\n"
                f"الخطأ: ({upload_error})\n"
            )
            if artifact_saved and artifact_filename:
                final_message += f"\nالملف متاح كـ Artifact باسم: {artifact_filename}"


        return {
            "output_filename": out_name,
            "slides_generated": topics_count,
            "gcs_bucket": bucket or None,
            "gcs_object": blob_name if bucket else None,
            "gcs_url": gcs_url,
            "signed_url": signed_url,
            "upload_error": upload_error,
            "artifact_saved": artifact_saved,
            "artifact_filename": artifact_filename,
            "final_message": final_message,
        }
